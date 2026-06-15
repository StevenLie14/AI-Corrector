import asyncio
import base64
from io import BytesIO

import fitz
from fastapi import APIRouter, File, UploadFile
from fastapi.responses import HTMLResponse
from pptx import Presentation

from schemas import DebugExtractResponse, DebugImagesResponse
from utils import chunk_text, extract_text

router = APIRouter(tags=["Debug"])

_IMAGE_MIN_BYTES = 10240


def _extract_images(file_bytes: bytes, filename: str) -> list[dict]:
    filename_lower = filename.lower()
    images = []

    if filename_lower.endswith(".pdf"):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(len(doc)):
            page = doc[page_num]
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                images.append({
                    "source": f"PDF page {page_num + 1}",
                    "size_bytes": len(image_bytes),
                    "skipped": len(image_bytes) < _IMAGE_MIN_BYTES,
                    "data": base64.b64encode(image_bytes).decode("utf-8") if len(image_bytes) >= _IMAGE_MIN_BYTES else None,
                })

    elif filename_lower.endswith(".pptx"):
        prs = Presentation(BytesIO(file_bytes))
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, "image"):
                    continue
                blob = shape.image.blob
                images.append({
                    "source": f"PPTX slide {slide_num + 1}",
                    "size_bytes": len(blob),
                    "skipped": len(blob) < _IMAGE_MIN_BYTES,
                    "data": base64.b64encode(blob).decode("utf-8") if len(blob) >= _IMAGE_MIN_BYTES else None,
                })

    return images


@router.post(
    "/debug/extract",
    response_model=DebugExtractResponse,
    summary="Extract text and chunks from a document",
    description=(
        "Parse a document and return its extracted raw text, text chunks, and vision token usage. "
        "Useful for verifying that text extraction and chunking work correctly before feeding a file."
    ),
)
async def debug_extract(file: UploadFile = File(..., description="PDF, PPT, PPTX, DOCX, or TXT file to inspect")):
    file_bytes = await file.read()
    raw_text, vision_tokens = await asyncio.to_thread(extract_text, file_bytes, file.filename)
    chunks = chunk_text(raw_text)
    return {
        "filename": file.filename,
        "vision_tokens_used": vision_tokens,
        "total_words": len(raw_text.split()),
        "total_chunks": len(chunks),
        "raw_text": raw_text,
        "chunks": chunks,
    }


@router.post(
    "/debug/images",
    response_model=DebugImagesResponse,
    summary="List images found in a document",
    description=(
        "Extract and list all images found in a PDF or PPTX file. "
        "Images smaller than 10 KB are flagged as `skipped`. "
        "Base64-encoded image data is included for images that pass the size threshold."
    ),
)
async def debug_images(file: UploadFile = File(..., description="PDF or PPTX file to inspect")):
    file_bytes = await file.read()
    images = await asyncio.to_thread(_extract_images, file_bytes, file.filename)
    return {
        "filename": file.filename,
        "total_images_found": len(images),
        "total_processed": sum(1 for img in images if not img["skipped"]),
        "total_skipped": sum(1 for img in images if img["skipped"]),
        "images": images,
    }


@router.post(
    "/debug/images/view",
    response_class=HTMLResponse,
    summary="View images from a document in a browser",
    description=(
        "Returns an HTML page that renders all images found in a PDF or PPTX file as visual cards. "
        "Images smaller than 10 KB are shown as greyed-out skipped cards."
    ),
    responses={200: {"content": {"text/html": {}}, "description": "HTML image viewer page"}},
)
async def debug_images_view(file: UploadFile = File(..., description="PDF or PPTX file to inspect")):
    file_bytes = await file.read()
    images = await asyncio.to_thread(_extract_images, file_bytes, file.filename)

    cards = ""
    for img in images:
        label = f"{img['source']} — {img['size_bytes']:,} bytes"
        if img["skipped"]:
            cards += f'<div class="card skipped"><p>{label} <span class="badge">SKIPPED (too small)</span></p></div>'
        else:
            cards += (
                f'<div class="card">'
                f'<p>{label}</p>'
                f'<img src="data:image/jpeg;base64,{img["data"]}" />'
                f'</div>'
            )

    html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>Debug Images — {file.filename}</title>
<style>
  body {{ font-family: sans-serif; padding: 20px; background: #f5f5f5; }}
  h1 {{ font-size: 18px; margin-bottom: 4px; }}
  .meta {{ color: #666; font-size: 14px; margin-bottom: 20px; }}
  .card {{ background: white; border-radius: 8px; padding: 12px; margin-bottom: 16px; box-shadow: 0 1px 4px rgba(0,0,0,.1); }}
  .card img {{ max-width: 100%; border: 1px solid #ddd; border-radius: 4px; display: block; margin-top: 8px; }}
  .card p {{ margin: 0; font-size: 13px; color: #333; }}
  .skipped {{ opacity: .5; }}
  .badge {{ background: #e0e0e0; padding: 2px 6px; border-radius: 4px; font-size: 11px; }}
</style>
</head>
<body>
<h1>Images in: {file.filename}</h1>
<p class="meta">
  Found {len(images)} image(s) —
  {sum(1 for i in images if not i["skipped"])} processed,
  {sum(1 for i in images if i["skipped"])} skipped
</p>
{cards}
</body>
</html>"""
    return HTMLResponse(content=html)
