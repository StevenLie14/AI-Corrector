import hashlib
import os
import re
from concurrent.futures import ThreadPoolExecutor
from html.parser import HTMLParser
from io import BytesIO

from config.constants import CHUNK_SIZE, CHUNK_OVERLAP

import fitz

from pptx import Presentation

from .image import get_image_description

_SKIP_TAGS = {"script", "style", "head", "nav", "footer", "header", "iframe", "noscript"}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._skip_depth = 0
        self._parts: list[str] = []

    def handle_starttag(self, tag, attrs):
        if tag in _SKIP_TAGS:
            self._skip_depth += 1

    def handle_endtag(self, tag):
        if tag in _SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data):
        if self._skip_depth == 0:
            stripped = data.strip()
            if stripped:
                self._parts.append(stripped)

    def get_text(self) -> str:
        return "\n".join(self._parts)


def extract_html_text(html_bytes: bytes) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(html_bytes.decode("utf-8", errors="ignore"))
    raw = parser.get_text()
    return re.sub(r"\n{3,}", "\n\n", raw).strip()

def extract_text(file_bytes: bytes, filename: str, is_student_answer: bool = False) -> tuple[str, int]:
    try:
        return _parse_file(BytesIO(file_bytes), file_bytes, filename, is_student_answer)
    except Exception as e:
        raise Exception(f"Error extracting text: {str(e)}")


def _parse_file(file_stream, file_bytes: bytes, filename: str, is_student_answer: bool = False) -> tuple[str, int]:
    filename_lower = filename.lower()
    images_to_process = []

    if filename_lower.endswith(".txt"):
        return file_bytes.decode("utf-8", errors="ignore"), 0

    elif filename_lower.endswith(".docx"):
        import docx
        _BLIP = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        _EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
        doc = docx.Document(BytesIO(file_bytes))
        img_rels = {}
        for rel_id, rel in doc.part.rels.items():
            if "image" in rel.reltype:
                try:
                    img_rels[rel_id] = rel.target_part.blob
                except Exception:
                    pass
        full_text = ""
        seen_blip_ids: set = set()
        for para in doc.paragraphs:
            para_content = "".join(run.text for run in para.runs)
            for blip in para._p.findall(".//" + _BLIP):
                r_embed = blip.get(_EMBED)
                if r_embed and r_embed in img_rels and r_embed not in seen_blip_ids:
                    seen_blip_ids.add(r_embed)
                    placeholder = f"[IMAGE_PLACEHOLDER_{len(images_to_process)}]"
                    para_content += "\n" + placeholder + "\n"
                    images_to_process.append((img_rels[r_embed], "DOCX"))
            if para_content.strip():
                full_text += para_content + "\n"

    elif filename_lower.endswith(".pdf"):
        full_text = ""
        seen_xrefs = set()
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            blocks = []
            for block in page.get_text("blocks"):
                if block[6] == 0:
                    blocks.append(("text", block[1], block[4]))

            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                y0 = rects[0].y0
                clip = rects[0]
                mat = fitz.Matrix(2, 2)
                pm = page.get_pixmap(matrix=mat, clip=clip)
                img_bytes = pm.tobytes("png")
                placeholder = f"[IMAGE_PLACEHOLDER_{len(images_to_process)}]"
                blocks.append(("image", y0, placeholder, img_bytes))
                images_to_process.append((img_bytes, "PDF"))

            blocks.sort(key=lambda b: b[1])
            for block in blocks:
                if block[0] == "text":
                    full_text += block[2].strip() + "\n"
                else:
                    full_text += block[2] + "\n"

            full_text += "\n"

    elif filename_lower.endswith(".pptx"):
        full_text = ""
        prs = Presentation(file_stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"
                if hasattr(shape, "image"):
                    placeholder = f"[IMAGE_PLACEHOLDER_{len(images_to_process)}]"
                    full_text += placeholder + "\n"
                    images_to_process.append((shape.image.blob, "PPTX"))

    elif filename_lower.endswith(".ppt"):
        import tempfile
        import ppt2txt

        with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        try:
            parsed = ppt2txt.process(tmp_path)
            content_dict = parsed.get("content", {})
            sorted_keys = sorted(content_dict.keys(), key=lambda k: int(k) if k.isdigit() else k)
            full_text = "\n".join(content_dict[k] for k in sorted_keys) + "\n"
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        return full_text, 0

    elif filename_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")):
        full_text = "[IMAGE_PLACEHOLDER_0]\n"
        images_to_process.append((file_bytes, "IMAGE"))

    else:
        raise ValueError("Unsupported file type. Supported: PDF, PPT, PPTX, TXT, DOCX, PNG, JPG, JPEG, GIF, BMP, WEBP")

    if images_to_process:
        full_text, vision_tokens = _replace_image_placeholders(full_text, images_to_process, is_student_answer)
        return full_text, vision_tokens

    return full_text, 0


def _replace_image_placeholders(
    full_text: str,
    images_to_process: list,
    is_student_answer: bool = False,
    cache: dict | None = None,
    cache_keys: list | None = None,
    strict: bool = False,
) -> tuple[str, int]:
    """Ganti placeholder gambar dengan deskripsi dari vision model.

    `cache` + `cache_keys` (opsional) memungkinkan dedup lintas pemanggilan: gambar
    yang identik (mis. logo yang muncul di tiap halaman PDF) cukup dideskripsikan
    sekali, sisanya memakai hasil yang sama. Tanpa ini, dokumen 151 halaman dengan
    1 logo berulang menghabiskan 147 panggilan vision untuk 5 gambar unik.

    `strict` menentukan apa yang terjadi kalau vision gagal setelah habis percobaan ulang:

    - `True` (jalur feed): lempar exception. Materi yang masuk index tanpa isi gambarnya
      adalah kerusakan senyap - dilaporkan sukses, tidak ada yang tahu isinya bolong.
      Lebih baik gagal terang-terangan; reconciler akan mengantre ulang di sapuan berikutnya.
    - `False` (jalur penilaian jawaban mahasiswa): pertahankan perilaku lama, lewati gambarnya
      saja. Permintaan penilaian bersifat interaktif dan tidak punya mekanisme antre-ulang,
      jadi menggagalkan seluruh permintaan karena satu gambar justru lebih merugikan.
    """
    replacements = {}
    tasks = []

    for i, (img_bytes, source) in enumerate(images_to_process):
        placeholder = f"[IMAGE_PLACEHOLDER_{i}]"
        idx = full_text.find(placeholder)

        if idx == -1:
            replacements[placeholder] = ""
            continue

        key = cache_keys[i] if cache_keys is not None and i < len(cache_keys) else None
        if cache is not None and key is not None and key in cache:
            replacements[placeholder] = cache[key]
            continue

        start = max(0, idx - 2000)
        end = min(len(full_text), idx + len(placeholder) + 2000)
        context = full_text[start:end].replace(placeholder, "\n[GAMBAR INI]\n")
        tasks.append((placeholder, img_bytes, context, source, key))

    total_vision_tokens = 0

    if tasks:
        def describe(task_info):
            ph, img_bytes, ctx, src, k = task_info
            try:
                desc, tokens = get_image_description(img_bytes, context_text=ctx, is_student_answer=is_student_answer)
                if not desc or (desc.strip().upper() == "SKIP" and not is_student_answer):
                    return ph, "", tokens, k
                return ph, f"\n[Deskripsi Gambar: {desc}]\n", tokens, k
            except Exception as e:
                # JANGAN diubah jadi hasil kosong diam-diam pada jalur feed: kosong berarti
                # "SKIP" (dekoratif) di baris di atas, jadi kegagalan tidak akan bisa dibedakan
                # dari pelewatan yang disengaja.
                if strict:
                    raise RuntimeError(f"Gagal mendeskripsikan gambar {ph} dari {src}: {e}") from e
                print(f"Gagal mengekstrak gambar {ph} dari {src}: {e}")
                return ph, "", 0, None

        with ThreadPoolExecutor(max_workers=10) as executor:
            for ph, replacement, tokens, k in executor.map(describe, tasks):
                replacements[ph] = replacement
                total_vision_tokens += tokens
                if cache is not None and k is not None:
                    cache[k] = replacement

    for placeholder, replacement in replacements.items():
        full_text = full_text.replace(placeholder, replacement)

    return full_text, total_vision_tokens


def extract_pages(file_bytes: bytes, filename: str) -> list[tuple[int, str, int]]:
    """Returns list of (page_num, text, vision_tokens) per page/slide."""
    try:
        return _parse_file_by_pages(BytesIO(file_bytes), file_bytes, filename)
    except Exception as e:
        raise Exception(f"Error extracting pages: {str(e)}")


def _parse_file_by_pages(file_stream, file_bytes: bytes, filename: str) -> list[tuple[int, str, int]]:
    filename_lower = filename.lower()

    if filename_lower.endswith(".txt"):
        return [(1, file_bytes.decode("utf-8", errors="ignore"), 0)]

    elif filename_lower.endswith(".docx"):
        import docx
        _BLIP = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        _EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
        doc = docx.Document(BytesIO(file_bytes))
        img_rels = {}
        for rel_id, rel in doc.part.rels.items():
            if "image" in rel.reltype:
                try:
                    img_rels[rel_id] = rel.target_part.blob
                except Exception:
                    pass
        images_to_process = []
        seen_blip_ids: set = set()
        full_text = ""
        for para in doc.paragraphs:
            para_content = "".join(run.text for run in para.runs)
            for blip in para._p.findall(".//" + _BLIP):
                r_embed = blip.get(_EMBED)
                if r_embed and r_embed in img_rels and r_embed not in seen_blip_ids:
                    seen_blip_ids.add(r_embed)
                    placeholder = f"[IMAGE_PLACEHOLDER_{len(images_to_process)}]"
                    para_content += "\n" + placeholder + "\n"
                    images_to_process.append((img_rels[r_embed], "DOCX"))
            if para_content.strip():
                full_text += para_content + "\n"
        if images_to_process:
            full_text, vtokens = _replace_image_placeholders(full_text, images_to_process, strict=True)
        else:
            vtokens = 0
        return [(1, full_text.strip(), vtokens)]

    elif filename_lower.endswith(".pptx"):
        pages = []
        prs = Presentation(file_stream)
        # Cache deskripsi per gambar, berlaku untuk SELURUH deck. Logo/elemen template
        # yang muncul di banyak slide cukup sekali dikirim ke vision model.
        deck_image_cache: dict = {}
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""
            slide_images = []
            slide_image_keys = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
                if hasattr(shape, "image"):
                    placeholder = f"[IMAGE_PLACEHOLDER_{len(slide_images)}]"
                    slide_text += placeholder + "\n"
                    image = shape.image
                    # PPTX tidak punya xref seperti PDF; pakai hash isi gambar sbg kunci.
                    key = getattr(image, "sha1", None)
                    if key is None:
                        key = hashlib.sha1(image.blob).hexdigest()
                    # gambar yang sudah pernah dideskripsikan tidak perlu dibaca ulang
                    blob = b"" if key in deck_image_cache else image.blob
                    slide_images.append((blob, "PPTX"))
                    slide_image_keys.append(key)
            if slide_images:
                slide_text, vtokens = _replace_image_placeholders(
                    slide_text, slide_images, cache=deck_image_cache, cache_keys=slide_image_keys,
                    strict=True
                )
            else:
                vtokens = 0
            pages.append((slide_num, slide_text.strip(), vtokens))
        return pages

    elif filename_lower.endswith(".pdf"):
        pages = []
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        # Cache deskripsi per xref, berlaku untuk SELURUH dokumen. Gambar yang sama
        # (logo/header berulang) cukup sekali dikirim ke vision model.
        doc_image_cache: dict = {}
        for page_num, page in enumerate(doc, 1):
            page_text = ""
            page_images = []
            page_image_keys = []
            blocks = []
            seen_xrefs = set()
            for block in page.get_text("blocks"):
                if block[6] == 0:
                    blocks.append(("text", block[1], block[4]))
            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                y0 = rects[0].y0
                clip = rects[0]
                placeholder = f"[IMAGE_PLACEHOLDER_{len(page_images)}]"
                if xref in doc_image_cache:
                    # sudah pernah dideskripsikan: lewati render pixmap dan panggilan vision
                    img_bytes = b""
                else:
                    mat = fitz.Matrix(2, 2)
                    pm = page.get_pixmap(matrix=mat, clip=clip)
                    img_bytes = pm.tobytes("png")
                blocks.append(("image", y0, placeholder, img_bytes))
                page_images.append((img_bytes, "PDF"))
                page_image_keys.append(xref)
            blocks.sort(key=lambda b: b[1])
            for block in blocks:
                page_text += (block[2].strip() if block[0] == "text" else block[2]) + "\n"
            if page_images:
                page_text, vtokens = _replace_image_placeholders(
                    page_text, page_images, cache=doc_image_cache, cache_keys=page_image_keys,
                    strict=True
                )
            else:
                vtokens = 0
            pages.append((page_num, page_text.strip(), vtokens))
        return pages

    elif filename_lower.endswith(".ppt"):
        import tempfile
        import ppt2txt
        with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            parsed = ppt2txt.process(tmp_path)
            content_dict = parsed.get("content", {})
            sorted_keys = sorted(content_dict.keys(), key=lambda k: int(k) if k.isdigit() else k)
            return [(int(k) if k.isdigit() else i + 1, content_dict[k], 0) for i, k in enumerate(sorted_keys)]
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    else:
        raise ValueError("Unsupported file type. Supported: PDF, PPT, PPTX, TXT, DOCX")


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list:
    words = text.split()
    step = chunk_size - overlap
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), step)]
