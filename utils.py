from fastapi import UploadFile
from pypdf import PdfReader
from pptx import Presentation
from config import embed_client
from io import BytesIO
import os
import base64
import httpx

def get_image_description(image_bytes: bytes, context_text: str = "") -> str:
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    url = os.getenv("MULTI_MODAL_URL")
    api_key = os.getenv("MULTI_MODAL_KEY")
    
    if not url or not api_key:
        return ""
        
    headers = {
        "Content-Type": "application/json",
        "api-key": api_key
    }
    
    prompt_text = "Tolong ekstrak dan deskripsikan informasi tekstual serta konteks visual penting (seperti grafik, diagram) dari gambar ini. Hanya berikan teks yang diekstrak dan deskripsinya, jangan berhalusinasi."
    if context_text:
        prompt_text += f"\n\nBerikut adalah teks konteks dari dokumen di sekitar gambar ini (sebelum dan sesudah gambar, ditandai dengan [GAMBAR INI], berguna jika gambar merujuk pada 'slide sebelumnya' atau 'slide selanjutnya'):\n{context_text}"
    
    payload = {
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": prompt_text
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}"
                        }
                    }
                ]
            }
        ],
        "max_tokens": 1000
    }
    
    try:
        with httpx.Client() as client:
            response = client.post(url, headers=headers, json=payload, timeout=60.0)
            if response.status_code == 200:
                return response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        print(f"Error calling multi-modal API: {e}")
    
    return ""

def extract_text(file_input, filename: str = None) -> str:
    text = ""
    try:
        if isinstance(file_input, UploadFile):
            file_stream = file_input.file
            filename_to_check = file_input.filename
        elif isinstance(file_input, bytes):
            if not filename:
                raise ValueError("Filename is required when passing bytes")
            file_stream = BytesIO(file_input)
            filename_to_check = filename
        else:
            raise ValueError("Input must be UploadFile or bytes")
        
        filename_lower = filename_to_check.lower()
        full_text = ""
        images_to_process = []

        if filename_lower.endswith(".pdf"):
            reader = PdfReader(file_stream)
            for page in reader.pages:
                full_text += page.extract_text() + "\n"
                if hasattr(page, 'images'):
                    for image_file_object in page.images:
                        if len(image_file_object.data) < 10240:
                            continue
                        placeholder = f"[IMAGE_PLACEHOLDER_{len(images_to_process)}]"
                        full_text += placeholder + "\n"
                        images_to_process.append((image_file_object.data, "PDF"))
        elif filename_lower.endswith(".pptx"):
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
                    if hasattr(shape, "image"):
                        if len(shape.image.blob) < 10240:
                            continue
                        placeholder = f"[IMAGE_PLACEHOLDER_{len(images_to_process)}]"
                        full_text += placeholder + "\n"
                        images_to_process.append((shape.image.blob, "PPTX"))
        elif filename_lower.endswith(".ppt"):
            import tempfile
            import ppt2txt
            
            if isinstance(file_input, bytes):
                file_bytes = file_input
            else:
                file_stream.seek(0)
                file_bytes = file_stream.read()
                file_stream.seek(0)
                
            with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
                
            try:
                parsed = ppt2txt.process(temp_file_path)
                content_dict = parsed.get("content", {})
                sorted_keys = sorted(content_dict.keys(), key=lambda k: int(k) if k.isdigit() else k)
                full_text = "\n".join(content_dict[k] for k in sorted_keys) + "\n"
            finally:
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)
        else:
            raise ValueError("Only PDF, PPT, and PPTX files are supported")
            
        if images_to_process:
            from concurrent.futures import ThreadPoolExecutor
            
            placeholders_to_replace = {}
            tasks = []
            
            for i, (img_bytes, source) in enumerate(images_to_process):
                placeholder = f"[IMAGE_PLACEHOLDER_{i}]"
                idx = full_text.find(placeholder)
                
                if idx != -1:
                    start_idx = max(0, idx - 2000)
                    end_idx = min(len(full_text), idx + len(placeholder) + 2000)
                    context_text = full_text[start_idx:end_idx].replace(placeholder, "\n[GAMBAR INI]\n")
                    tasks.append((placeholder, img_bytes, context_text, source))
                else:
                    placeholders_to_replace[placeholder] = ""
            
            if tasks:
                def desc_task(task_info):
                    ph, bytes_data, ctx, src = task_info
                    try:
                        desc = get_image_description(bytes_data, context_text=ctx)
                        return ph, f"\n[Deskripsi Gambar: {desc}]\n" if desc else ""
                    except Exception as e:
                        print(f"Gagal mengekstrak gambar {ph} dari {src}: {e}")
                        return ph, ""
                
                with ThreadPoolExecutor(max_workers=10) as executor:
                    results = list(executor.map(desc_task, tasks))
                    for ph, replacement in results:
                        placeholders_to_replace[ph] = replacement
            
            for placeholder, replacement in placeholders_to_replace.items():
                full_text = full_text.replace(placeholder, replacement)
                    
        text = full_text
        
    except Exception as e:
        raise Exception(f"Error extracting text: {str(e)}")
    
    return text

def chunk_text(text: str, chunk_size: int = 400) -> list:
    words = text.split()
    chunks = [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
    return chunks

def get_embedding(text: str) -> list:
    response = embed_client.embeddings.create(
        input=text,
        model="text-embedding-3-small"
    )
    return response.data[0].embedding